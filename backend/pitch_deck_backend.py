#!/usr/bin/env python3
"""
AI-Powered Pitch Deck Builder Backend
Complete Flask application with AI integration for pitch deck generation
"""

import os
import json
import io
import logging
import re
from datetime import datetime
from typing import Dict, List, Optional, Any
from dataclasses import dataclass, asdict
from pathlib import Path

from flask import Flask, request, jsonify, send_file, render_template_string
from flask_cors import CORS
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from dotenv import load_dotenv

# For AI providers
try:
    import openai
except ImportError:
    openai = None

try:
    import anthropic
except ImportError:
    anthropic = None

# For file processing
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
    print("Warning: PyPDF2 not installed. PDF extraction will be limited.")

try:
    from docx import Document
except ImportError:
    Document = None
    print("Warning: python-docx not installed. DOCX extraction will be limited.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("Warning: python-pptx not installed. PowerPoint export will be limited.")

# Load environment variables
load_dotenv()

# Initialize Flask app
app = Flask(__name__)
CORS(app, origins=["http://localhost:*", "http://127.0.0.1:*", "https://*.onrender.com"])

# Set up rate limiting
limiter = Limiter(
    app=app,
    key_func=get_remote_address,
    default_limits=["100 per hour"],
    storage_uri="memory://"
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Configuration
class Config:
    """Application configuration"""
    SECRET_KEY = os.getenv('SECRET_KEY', 'dev-secret-key-change-in-production')
    AI_PROVIDER = os.getenv('AI_PROVIDER', 'openai')
    OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
    OPENAI_MODEL = os.getenv('OPENAI_MODEL', 'gpt-5')  # Default to GPT-5
    ANTHROPIC_API_KEY = os.getenv('ANTHROPIC_API_KEY')
    MAX_CONTENT_LENGTH = 50 * 1024 * 1024  # 50MB max file size
    USE_BUDGET_MODEL = os.getenv('USE_BUDGET_MODEL', 'false').lower() == 'true'

app.config.from_object(Config)

# Initialize AI clients
if Config.AI_PROVIDER == 'openai' and openai and Config.OPENAI_API_KEY:
    from openai import OpenAI
    ai_client = OpenAI(api_key=Config.OPENAI_API_KEY)
elif Config.AI_PROVIDER == 'anthropic' and anthropic and Config.ANTHROPIC_API_KEY:
    ai_client = anthropic.Anthropic(api_key=Config.ANTHROPIC_API_KEY)
else:
    ai_client = None
    logger.warning("No AI provider configured. Using mock responses.")

# Data models
@dataclass
class PitchSlide:
    """Individual pitch slide structure"""
    type: str
    title: str
    content: str
    order: int
    metadata: Optional[Dict] = None

@dataclass
class PitchDeck:
    """Complete pitch deck structure"""
    company_name: str
    tagline: str
    industry: str
    funding_stage: str
    slides: List[PitchSlide]
    contact_info: Optional[Dict] = None

@dataclass
class PitchRequest:
    """Pitch generation request"""
    company_name: str
    industry: str
    problem: str
    solution: str
    funding_stage: str = "seed"
    target_investors: Optional[str] = None
    team_size: Optional[int] = None
    current_traction: Optional[str] = None

# AI Pitch Generator
class AIPitchGenerator:
    """Handles AI-powered pitch deck generation"""
    
    @staticmethod
    def generate_with_openai(request: PitchRequest) -> Dict[str, Any]:
        """Generate pitch deck using OpenAI GPT-5 with full capabilities"""
        if not ai_client:
            return AIPitchGenerator.generate_mock_pitch(request)
        
        # Choose model based on configuration
        if Config.USE_BUDGET_MODEL:
            model = "gpt-3.5-turbo"
            max_tokens = 2000
            logger.info("Using GPT-3.5-turbo (budget mode)")
        else:
            model = "gpt-5-2025-08-07"  # Latest GPT-5 model
            max_tokens = None  # No limit - let GPT-5 use its full 128k output capacity
            logger.info(f"Using {model} with unlimited output")
        
        # Enhanced prompt for GPT-5's superior reasoning capabilities
        prompt = f"""
        You are a world-class pitch deck consultant who has helped startups raise over $10B in funding.
        Use your advanced reasoning capabilities to create the MOST COMPREHENSIVE, INVESTOR-READY pitch deck possible.
        
        Company Details:
        - Company Name: {request.company_name}
        - Industry: {request.industry}
        - Funding Stage: {request.funding_stage}
        - Team Size: {request.team_size or 'Not specified'}
        - Current Traction: {request.current_traction or 'Early stage'}
        - Target Investors: {request.target_investors or 'Tier 1 VCs'}
        
        Core Value Proposition:
        PROBLEM: {request.problem}
        SOLUTION: {request.solution}
        
        REASONING INSTRUCTIONS:
        1. First, analyze the market dynamics and competitive landscape
        2. Calculate realistic financial projections based on industry benchmarks
        3. Develop a compelling narrative arc throughout the deck
        4. Ensure each slide builds on the previous one
        5. Include specific metrics, data points, and evidence
        6. Create content that addresses likely investor concerns preemptively
        
        Generate an EXHAUSTIVE pitch deck with these slides:
        
        1. TITLE SLIDE: 
           - Company name, tagline, logo description
           - Contact details, website, location
           - Founding date and current stage
        
        2. PROBLEM (make it visceral and urgent):
           - Primary problem with specific cost to businesses/consumers
           - Secondary problems that compound the issue
           - Real customer quotes demonstrating pain (create realistic examples)
           - Market inefficiencies currently being exploited
           - Why this problem exists and persists
           - Consequences of not solving it
        
        3. SOLUTION (show breakthrough innovation):
           - Core solution with technical differentiation
           - How it works (detailed but accessible explanation)
           - Key features and benefits matrix
           - Technology stack and proprietary elements
           - Implementation timeline for customers
           - ROI calculation for customers
        
        4. MARKET OPPORTUNITY (be extremely detailed):
           - TAM: Global market with breakdown by region (NA, EU, APAC, LATAM, MEA)
           - TAM calculation methodology with sources
           - SAM: Addressable segments with rationale
           - SOM: 5-year capture plan with specific milestones
           - Market growth drivers (regulatory, technological, social)
           - Market timing analysis - why now?
           - Adjacent markets for expansion
        
        5. BUSINESS MODEL (show path to $1B):
           - Revenue streams with pricing for each tier
           - Customer segmentation and pricing strategy
           - Sales model (self-serve, inside sales, enterprise)
           - Customer acquisition strategy and channels
           - Unit economics deep dive:
             * CAC by channel
             * LTV by segment  
             * Gross margins
             * Contribution margins
             * Payback period
           - Network effects and virality coefficients
           - Expansion revenue strategy
        
        6. TRACTION & VALIDATION (prove product-market fit):
           - Monthly metrics for last 12 months (or projected)
           - Customer growth curve
           - Revenue growth with MRR/ARR progression
           - Key metrics: NPS, DAU/MAU, retention, churn
           - Customer logos and case studies
           - Testimonials with specific results
           - Product development milestones
           - Awards and recognition
           - Media coverage and PR wins
        
        7. COMPETITION & MOAT (show defensibility):
           - Competitive matrix (10+ competitors)
           - Feature comparison table
           - Pricing comparison
           - Your unfair advantages (list 5+)
           - Barriers to entry you're creating
           - Network effects and lock-in
           - IP and patents (filed or planned)
           - Why competitors can't copy you
           - M&A opportunities
        
        8. TEAM (prove you can execute):
           - Founders: Background, achievements, why them
           - Key employees with notable backgrounds
           - Advisory board with their contributions
           - Board members (if any)
           - Key hires planned next 12 months
           - Culture and values
           - Why this team wins
        
        9. FINANCIAL PROJECTIONS (be aggressive but defensible):
           - 5-year P&L with monthly detail for Year 1
           - Revenue build (bottom-up model)
           - Cost structure breakdown
           - Hiring plan by department
           - Key assumptions clearly stated
           - Sensitivity analysis
           - Break-even timeline
           - Cash flow projections
           - Scenario planning (base, upside, downside)
        
        10. THE ASK (be specific and ambitious):
           - Funding amount with exact figure
           - Valuation expectations (pre/post money)
           - Use of funds with percentages:
             * Product development
             * Sales & Marketing  
             * Operations
             * Team expansion
             * Working capital
           - 18-month milestone plan with KPIs
           - Expected next round timing and size
           - Strategic investors desired
           - Exit strategy and comparables
        
        11. APPENDIX (bonus comprehensive content):
           - Technical architecture diagram
           - Go-to-market playbook
           - Customer journey map
           - Partnership pipeline
           - PR and marketing strategy
           - Risk matrix and mitigation
           - Detailed competitor analysis
           - Customer testimonial videos (scripts)
           - Demo screenshots and flows
           - Press quotes and coverage
        
        Make every slide EXCEPTIONAL with:
        - Specific numbers (not ranges)
        - Real company names as examples
        - Industry-specific terminology
        - Compelling data visualizations described
        - Emotional hooks and storytelling
        - Clear calls to action
        
        Use your FULL REASONING CAPABILITIES to make this the best pitch deck ever created.
        Generate AS MUCH DETAIL AS POSSIBLE - use the full 128k output token capacity if needed.
        
        Return as comprehensive JSON with all content.
        """
        
        try:
            response = ai_client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system", 
                        "content": """You are the world's best pitch deck consultant with experience raising billions.
                        Use your advanced reasoning to create comprehensive, detailed, investor-ready content.
                        Do not hold back on detail - use as many tokens as needed to be thorough."""
                    },
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=max_tokens,  # None for GPT-5 = unlimited
                response_format={"type": "json_object"}
            )
            
            content = json.loads(response.choices[0].message.content)
            
            # Log token usage for GPT-5
            if hasattr(response, 'usage'):
                logger.info(f"GPT-5 token usage - Input: {response.usage.prompt_tokens}, "
                          f"Output: {response.usage.completion_tokens}, "
                          f"Reasoning: {getattr(response.usage, 'reasoning_tokens', 'N/A')}")
            
            logger.info(f"Successfully generated comprehensive pitch deck with {model}")
            return AIPitchGenerator.format_pitch_response(content, request)
            
        except Exception as e:
            logger.error(f"OpenAI generation failed: {e}")
            # Fallback to GPT-3.5 if GPT-5 fails
            if model != "gpt-3.5-turbo":
                logger.info("Falling back to GPT-3.5-turbo")
                Config.USE_BUDGET_MODEL = True
                return AIPitchGenerator.generate_with_openai(request)
            return AIPitchGenerator.generate_mock_pitch(request)
    
    @staticmethod
    def generate_with_anthropic(request: PitchRequest) -> Dict[str, Any]:
        """Generate pitch deck using Anthropic Claude"""
        if not ai_client:
            return AIPitchGenerator.generate_mock_pitch(request)
        
        prompt = f"""
        Create a compelling investor pitch deck for {request.company_name}, a {request.funding_stage}-stage {request.industry} startup.
        
        Problem: {request.problem}
        Solution: {request.solution}
        
        Generate content for all 10 standard pitch deck slides. Return as JSON.
        """
        
        try:
            response = ai_client.messages.create(
                model="claude-3-opus-20240229",
                max_tokens=2000,
                messages=[{"role": "user", "content": prompt}]
            )
            
            content = json.loads(response.content[0].text)
            return AIPitchGenerator.format_pitch_response(content, request)
            
        except Exception as e:
            logger.error(f"Anthropic generation failed: {e}")
            return AIPitchGenerator.generate_mock_pitch(request)
    
    @staticmethod
    def generate_mock_pitch(request: PitchRequest) -> Dict[str, Any]:
        """Generate mock pitch deck for testing/fallback"""
        
        # Funding stage specifics
        funding_amounts = {
            "pre-seed": "$500K - $1M",
            "seed": "$2M - $5M",
            "series-a": "$10M - $15M",
            "series-b": "$20M - $50M",
            "series-c": "$50M+"
        }
        
        funding_amount = funding_amounts.get(request.funding_stage, "$5M")
        
        # Industry-specific TAM
        tam_by_industry = {
            "fintech": "$1.5T",
            "healthcare": "$3.6T",
            "saas": "$200B",
            "ecommerce": "$5.5T",
            "edtech": "$350B",
            "ai": "$150B"
        }
        
        tam = tam_by_industry.get(request.industry.lower(), "$100B")
        
        return {
            "tagline": f"Revolutionizing {request.industry} through innovative technology",
            
            "title": f"""{request.company_name}

{f"Revolutionizing {request.industry} through innovative technology"}

Investor Deck | {request.funding_stage.title()} Round
{datetime.now().strftime('%B %Y')}""",
            
            "problem": f"""{request.problem}

Current Pain Points:
• Inefficient legacy systems costing businesses millions annually
• 73% of {request.industry} professionals report daily frustrations
• Average time wasted: 4.5 hours per week per employee
• Customer satisfaction scores declining 15% YoY
• No comprehensive solution exists in the market""",
            
            "solution": f"""{request.solution}

Our Approach:
• AI-powered automation reducing manual work by 80%
• Seamless integration with existing workflows
• Real-time analytics and insights dashboard
• Mobile-first design for modern teams
• Enterprise-grade security and compliance

Results: 10x faster, 50% cost reduction, 95% user satisfaction""",
            
            "market": f"""Market Opportunity

📊 Total Addressable Market (TAM): {tam}
   - Growing at 25% CAGR
   - Digital transformation driving demand

🎯 Serviceable Addressable Market (SAM): ${tam[1:-1]}B
   - Focus on mid-market and enterprise
   - {request.industry} segment specifically

🚀 Serviceable Obtainable Market (SOM): $500M
   - Realistic 5-year target
   - 1% market share achievable

Key Drivers:
• Regulatory changes forcing modernization
• Remote work acceleration
• Gen Z entering workforce""",
            
            "business_model": f"""Revenue Model

💰 Subscription (SaaS)
   • Starter: $99/month (freelancers)
   • Professional: $499/month (small teams)
   • Enterprise: $2,999/month (large orgs)
   • Custom pricing for 100+ seats

Additional Revenue Streams:
• Implementation services: $10-50K
• API access: Usage-based pricing
• Premium support: 20% of license fee
• Marketplace commissions: 15%

Unit Economics:
• CAC: $2,000 | LTV: $45,000 | LTV/CAC: 22.5x
• Gross Margin: 82%
• Payback Period: 8 months""",
            
            "traction": f"""Traction & Validation

📈 Growth Metrics:
• 10,000+ users across 500+ companies
• $2.5M ARR (growing 30% MoM)
• 120% net revenue retention
• NPS Score: 72
• 5-min average time to value

🏆 Key Achievements:
• Product Hunt #1 Product of the Day
• SOC 2 Type II certified
• 3 Fortune 500 pilots in progress
• Strategic partnership with Microsoft
• 50+ 5-star reviews on G2

📊 Usage Stats:
• 1M+ transactions processed monthly
• 99.99% uptime over last 12 months
• 3-minute average response time""",
            
            "competition": f"""Competitive Landscape

Direct Competitors:
🔴 Legacy Corp ($2B valuation)
   - Strength: Market share (35%)
   - Weakness: Outdated tech, poor UX
   - Our advantage: 10x faster, 50% cheaper

🟡 StartupX (Series B, $150M raised)
   - Strength: Strong marketing
   - Weakness: Limited features
   - Our advantage: Complete platform

🔵 BigTech's Solution
   - Strength: Brand recognition
   - Weakness: Not specialized
   - Our advantage: Industry focus

Competitive Advantages:
✅ Proprietary AI technology (3 patents pending)
✅ 5x faster implementation
✅ 50% lower TCO
✅ Best-in-class user experience
✅ Only solution with full mobile support""",
            
            "team": f"""Leadership Team

👤 CEO & Co-founder
• 10+ years in {request.industry}
• Previously VP at Fortune 500
• Scaled 2 startups to exit
• Harvard MBA

👤 CTO & Co-founder
• Ex-Google/Amazon engineer
• 15+ years building scalable systems
• PhD Computer Science, Stanford
• 20+ patents in AI/ML

👤 VP Sales
• Built $100M+ revenue teams
• Former CRO at unicorn startup
• Deep industry relationships

👤 VP Product
• Led product at 3 successful startups
• Human-centered design expert
• Previously at Apple

Advisory Board:
• Former CEO of [Industry Leader]
• Partner at Sequoia Capital
• Professor of AI at MIT""",
            
            "financials": f"""Financial Projections

Current Status:
• Monthly Burn: $250K
• Runway: 18 months
• Path to profitability: Q3 2025

3-Year Projections:
         Year 1    Year 2     Year 3
Revenue:  $5M      $22M       $65M
Gross:    $4.1M    $18M       $53M
EBITDA:   -$3M     $2M        $15M
Customers: 50      250        800

Key Assumptions:
• 15% monthly growth rate Year 1
• 80% gross margins maintained
• Sales efficiency improves 20% YoY
• Churn remains below 5% annually

Use of Funds:
• Product Development: 40%
• Sales & Marketing: 35%
• Operations: 15%
• General & Admin: 10%""",
            
            "ask": f"""The Ask

🎯 Raising: {funding_amount} {request.funding_stage.title()} Round

Use of Funds:
• Engineering: Hire 10 engineers to accelerate product roadmap
• Sales: Build enterprise sales team (8 reps)
• Marketing: Scale demand generation and brand
• Operations: Strengthen infrastructure for scale

Milestones (Next 18 Months):
✓ Launch AI-powered analytics suite (Q1)
✓ Expand to 3 new verticals (Q2)
✓ Achieve $10M ARR (Q3)
✓ Close 5 enterprise accounts (Q4)
✓ International expansion (Q1 Y2)
✓ Series {chr(ord(request.funding_stage[-1]) + 1).upper() if request.funding_stage.startswith('series') else 'A'} ready

Why Now:
• Market inflection point
• Proven product-market fit
• Team ready to scale
• Competition vulnerable

Contact: founders@{request.company_name.lower().replace(' ', '')}.com"""
        }
    
    @staticmethod
    def format_pitch_response(content: Dict, request: PitchRequest) -> Dict[str, Any]:
        """Format AI response into consistent structure"""
        return {
            "tagline": content.get("tagline", f"Transforming {request.industry}"),
            "title": content.get("title", f"{request.company_name} Pitch Deck"),
            "problem": content.get("problem", request.problem),
            "solution": content.get("solution", request.solution),
            "market": content.get("market", "Market opportunity analysis"),
            "business_model": content.get("business_model", "Subscription-based SaaS model"),
            "traction": content.get("traction", "Early traction and validation"),
            "competition": content.get("competition", "Competitive landscape"),
            "team": content.get("team", "Experienced leadership team"),
            "financials": content.get("financials", "Financial projections"),
            "ask": content.get("ask", f"Raising {request.funding_stage} round")
        }
    
    @classmethod
    def generate(cls, request: PitchRequest) -> Dict[str, Any]:
        """Main generation method"""
        if Config.AI_PROVIDER == 'openai':
            return cls.generate_with_openai(request)
        elif Config.AI_PROVIDER == 'anthropic':
            return cls.generate_with_anthropic(request)
        else:
            return cls.generate_mock_pitch(request)

# Content Extractor
class ContentExtractor:
    """Extract content from various file formats"""
    
    @staticmethod
    def extract_from_pdf(content: bytes) -> str:
        """Extract text from PDF"""
        if not PyPDF2:
            return ""
        
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            return ""
    
    @staticmethod
    def extract_from_docx(content: bytes) -> str:
        """Extract text from DOCX"""
        if not Document:
            return ""
        
        try:
            doc = Document(io.BytesIO(content))
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            return ""
    
    @staticmethod
    def extract_from_text(content: bytes) -> str:
        """Extract text from plain text file"""
        try:
            return content.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f"Text extraction error: {e}")
            return ""
    
    @classmethod
    def extract(cls, content: bytes, file_type: str) -> str:
        """Main extraction method"""
        if 'pdf' in file_type.lower():
            return cls.extract_from_pdf(content)
        elif 'docx' in file_type.lower() or 'word' in file_type.lower():
            return cls.extract_from_docx(content)
        else:
            return cls.extract_from_text(content)

# PowerPoint Exporter
class PowerPointExporter:
    """Export pitch deck to PowerPoint format"""
    
    @staticmethod
    def create_presentation(pitch_data: Dict) -> bytes:
        """Create PowerPoint presentation"""
        if not Presentation:
            logger.warning("python-pptx not installed")
            return None
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            
            # Set slide size to 16:9
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            slides_data = pitch_data.get('slides', [])
            
            for slide_data in slides_data:
                # Add slide with title and content layout
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                title = slide.shapes.title
                title.text = slide_data.get('title', 'Slide Title')
                
                # Set content
                content = slide.placeholders[1]
                content.text = slide_data.get('content', '')
                
                # Format text
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
            
            # Save to bytes
            ppt_bytes = io.BytesIO()
            prs.save(ppt_bytes)
            ppt_bytes.seek(0)
            
            return ppt_bytes.getvalue()
            
        except Exception as e:
            logger.error(f"PowerPoint export error: {e}")
            return None

# API Routes
@app.route('/')
def index():
    """Serve the pitch deck builder frontend"""
    frontend_path = Path(__file__).parent.parent / 'frontend' / 'pitch-deck-builder.html'
    
    try:
        if frontend_path.exists():
            with open(frontend_path, 'r', encoding='utf-8') as f:
                return f.read(), 200, {'Content-Type': 'text/html'}
    except Exception as e:
        logger.error(f"Error serving frontend: {e}")
    
    # Fallback HTML
    return """
    <!DOCTYPE html>
    <html>
    <head>
        <title>Pitch Deck Builder API</title>
        <style>
            body { 
                font-family: Arial, sans-serif; 
                max-width: 800px; 
                margin: 50px auto; 
                padding: 20px;
                background: #0a0a0a;
                color: #fff;
            }
            h1 { color: #ff6600; }
            .endpoint { 
                background: #1a1a1a; 
                padding: 15px; 
                margin: 10px 0; 
                border-radius: 5px;
                border-left: 4px solid #ff6600;
            }
            code { 
                background: #333; 
                padding: 2px 5px; 
                border-radius: 3px; 
            }
        </style>
    </head>
    <body>
        <h1>🚀 Pitch Deck Builder API</h1>
        <p>The API is running! Frontend file not found at expected location.</p>
        
        <h2>Available Endpoints:</h2>
        <div class="endpoint">
            <strong>POST</strong> <code>/api/generate-pitch</code> - Generate AI pitch deck
        </div>
        <div class="endpoint">
            <strong>POST</strong> <code>/api/extract-pitch</code> - Extract content from documents
        </div>
        <div class="endpoint">
            <strong>POST</strong> <code>/api/analyze-pitch</code> - Analyze existing pitch
        </div>
        <div class="endpoint">
            <strong>POST</strong> <code>/api/export/pptx</code> - Export as PowerPoint
        </div>
        <div class="endpoint">
            <strong>POST</strong> <code>/api/export/pdf</code> - Export as PDF
        </div>
        <div class="endpoint">
            <strong>GET</strong> <code>/health</code> - Health check
        </div>
    </body>
    </html>
    """, 200

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "ai_provider": Config.AI_PROVIDER,
        "ai_available": ai_client is not None,
        "pdf_support": PyPDF2 is not None,
        "docx_support": Document is not None,
        "pptx_support": Presentation is not None
    })

@app.route('/api/generate-pitch', methods=['POST'])
@limiter.limit("10 per hour")
def generate_pitch():
    """Generate complete pitch deck using AI"""
    try:
        data = request.json
        
        # Validate required fields
        if not data.get('company_name') or not data.get('problem') or not data.get('solution'):
            return jsonify({"error": "Company name, problem, and solution are required"}), 400
        
        # Create request object
        pitch_request = PitchRequest(
            company_name=data.get('company_name'),
            industry=data.get('industry', 'technology'),
            problem=data.get('problem'),
            solution=data.get('solution'),
            funding_stage=data.get('funding_stage', 'seed'),
            target_investors=data.get('target_investors'),
            team_size=data.get('team_size'),
            current_traction=data.get('current_traction')
        )
        
        # Generate pitch content
        pitch_content = AIPitchGenerator.generate(pitch_request)
        
        return jsonify(pitch_content)
    
    except Exception as e:
        logger.error(f"Pitch generation error: {e}")
        return jsonify({"error": "Failed to generate pitch deck"}), 500

@app.route('/api/extract-pitch', methods=['POST'])
@limiter.limit("10 per hour")
def extract_pitch_content():
    """Extract and structure content from uploaded documents"""
    try:
        data = request.json
        content = data.get('content', '')
        file_count = data.get('fileCount', 1)
        
        if not content:
            return jsonify({"error": "No content provided"}), 400
        
        # Use AI to extract structured pitch content
        extracted_data = extract_pitch_with_ai(content)
        
        return jsonify(extracted_data)
    
    except Exception as e:
        logger.error(f"Content extraction error: {e}")
        return jsonify({"error": "Failed to extract content"}), 500

def extract_pitch_with_ai(content: str) -> Dict[str, Any]:
    """Extract pitch deck content using AI"""
    
    if ai_client and Config.AI_PROVIDER == 'openai':
        try:
            prompt = f"""
            Extract pitch deck information from this content and structure it for investor slides:
            
            {content[:3000]}
            
            Return JSON with:
            - company: object with name, tagline, website, email
            - slides: array of objects with type, title, and content
            
            Identify and extract:
            - Company information
            - Problem statement
            - Solution description
            - Market size/opportunity
            - Business model
            - Traction metrics
            - Team information
            - Financial data
            - Funding requirements
            """
            
            response = ai_client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": "Extract and structure pitch deck content from documents."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=1500,
                response_format={"type": "json_object"}
            )
            
            return json.loads(response.choices[0].message.content)
            
        except Exception as e:
            logger.error(f"AI extraction failed: {e}")
    
    # Fallback: Basic extraction
    return {
        "company": {
            "name": extract_company_name(content),
            "tagline": "Extracted from documents",
            "website": extract_url(content),
            "email": extract_email(content)
        },
        "slides": [
            {"type": "custom", "title": "Extracted Content", "content": content[:1000]}
        ]
    }

def extract_company_name(content: str) -> str:
    """Extract company name from content"""
    lines = content.split('\n')
    for line in lines[:10]:  # Check first 10 lines
        if len(line) > 2 and len(line) < 50 and not any(c in line for c in ['@', 'http', '.']):
            return line.strip()
    return "Company Name"

def extract_url(content: str) -> str:
    """Extract URL from content"""
    url_pattern = r'https?://[^\s]+'
    match = re.search(url_pattern, content)
    return match.group() if match else ""

def extract_email(content: str) -> str:
    """Extract email from content"""
    email_pattern = r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
    match = re.search(email_pattern, content)
    return match.group() if match else ""

@app.route('/api/analyze-pitch', methods=['POST'])
@limiter.limit("5 per hour")
def analyze_pitch():
    """Analyze and score an existing pitch deck"""
    try:
        data = request.json
        slides = data.get('slides', [])
        
        if not slides:
            return jsonify({"error": "No slides provided"}), 400
        
        # Analyze pitch quality
        analysis = analyze_pitch_quality(slides)
        
        return jsonify(analysis)
    
    except Exception as e:
        logger.error(f"Pitch analysis error: {e}")
        return jsonify({"error": "Failed to analyze pitch"}), 500

def analyze_pitch_quality(slides: List[Dict]) -> Dict[str, Any]:
    """Analyze pitch deck quality and completeness"""
    
    # Check for essential slides
    slide_types = [slide.get('type', '') for slide in slides]
    
    essential_slides = ['problem', 'solution', 'market', 'business-model', 'team', 'ask']
    missing_slides = [s for s in essential_slides if s not in slide_types]
    
    # Calculate scores
    completeness_score = (len(essential_slides) - len(missing_slides)) / len(essential_slides) * 100
    
    # Analyze content quality
    total_content_length = sum(len(slide.get('content', '')) for slide in slides)
    avg_content_length = total_content_length / len(slides) if slides else 0
    
    quality_score = min(100, (avg_content_length / 200) * 100)  # Assume 200 chars is good
    
    # Generate recommendations
    recommendations = []
    
    if missing_slides:
        recommendations.append(f"Add missing slides: {', '.join(missing_slides)}")
    
    if avg_content_length < 100:
        recommendations.append("Add more detail to your slides")
    
    if 'traction' not in slide_types:
        recommendations.append("Include traction metrics to show validation")
    
    if 'financials' not in slide_types:
        recommendations.append("Add financial projections")
    
    # Overall score
    overall_score = (completeness_score + quality_score) / 2
    
    return {
        "overall_score": round(overall_score, 1),
        "completeness_score": round(completeness_score, 1),
        "quality_score": round(quality_score, 1),
        "slide_count": len(slides),
        "missing_slides": missing_slides,
        "recommendations": recommendations,
        "strengths": identify_strengths(slides),
        "verdict": "Ready to pitch!" if overall_score > 80 else "Needs improvement"
    }

def identify_strengths(slides: List[Dict]) -> List[str]:
    """Identify strengths in the pitch deck"""
    strengths = []
    
    slide_types = [slide.get('type', '') for slide in slides]
    
    if 'problem' in slide_types and 'solution' in slide_types:
        strengths.append("Clear problem-solution fit")
    
    if 'team' in slide_types:
        strengths.append("Team credentials included")
    
    if 'market' in slide_types:
        strengths.append("Market opportunity defined")
    
    if len(slides) >= 10:
        strengths.append("Comprehensive deck")
    
    return strengths

@app.route('/api/export/pptx', methods=['POST'])
@limiter.limit("20 per hour")
def export_powerpoint():
    """Export pitch deck as PowerPoint"""
    try:
        data = request.json
        
        if not data.get('slides'):
            return jsonify({"error": "No slides provided"}), 400
        
        # Create PowerPoint
        pptx_bytes = PowerPointExporter.create_presentation(data)
        
        if not pptx_bytes:
            return jsonify({"error": "PowerPoint export not available. Install python-pptx."}), 503
        
        return send_file(
            io.BytesIO(pptx_bytes),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f"pitch_deck_{datetime.now().strftime('%Y%m%d')}.pptx"
        )
    
    except Exception as e:
        logger.error(f"PowerPoint export error: {e}")
        return jsonify({"error": "Failed to export PowerPoint"}), 500

@app.route('/api/export/pdf', methods=['POST'])
@limiter.limit("20 per hour")
def export_pdf():
    """Export pitch deck as PDF"""
    try:
        data = request.json
        slides = data.get('slides', [])
        company_name = data.get('company', {}).get('name', 'Pitch Deck')
        
        # Generate HTML for PDF
        html_content = generate_pitch_html(slides, company_name)
        
        # Try to use WeasyPrint if available
        try:
            from weasyprint import HTML
            pdf = HTML(string=html_content).write_pdf()
            
            return send_file(
                io.BytesIO(pdf),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f"pitch_deck_{datetime.now().strftime('%Y%m%d')}.pdf"
            )
        except ImportError:
            # Fallback: return HTML for browser printing
            return html_content, 200, {'Content-Type': 'text/html'}
    
    except Exception as e:
        logger.error(f"PDF export error: {e}")
        return jsonify({"error": "Failed to export PDF"}), 500

def generate_pitch_html(slides: List[Dict], company_name: str) -> str:
    """Generate HTML for pitch deck"""
    
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <title>{company_name} - Pitch Deck</title>
        <style>
            @page {{ size: A4 landscape; margin: 0; }}
            body {{ 
                font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Arial, sans-serif;
                margin: 0;
                padding: 0;
                background: #0a0a0a;
                color: #ffffff;
            }}
            .slide {{
                width: 100%;
                height: 100vh;
                padding: 60px;
                box-sizing: border-box;
                display: flex;
                flex-direction: column;
                justify-content: center;
                page-break-after: always;
                background: linear-gradient(135deg, #0a0a0a 0%, #1a1a1a 100%);
            }}
            .slide h1 {{
                color: #ff6600;
                font-size: 48px;
                margin-bottom: 30px;
                text-align: center;
            }}
            .slide h2 {{
                color: #ff6600;
                font-size: 36px;
                margin-bottom: 30px;
            }}
            .slide p {{
                font-size: 20px;
                line-height: 1.8;
                margin-bottom: 20px;
            }}
            .slide ul {{
                font-size: 20px;
                line-height: 2;
                margin-left: 30px;
            }}
            .slide li {{
                margin-bottom: 15px;
            }}
            @media print {{
                .slide {{ height: 100vh; }}
            }}
        </style>
    </head>
    <body>
    """
    
    for slide in slides:
        content = slide.get('content', '').replace('\n', '<br>')
        
        # Convert bullet points
        if '•' in content:
            lines = content.split('<br>')
            formatted_lines = []
            in_list = False
            
            for line in lines:
                if line.strip().startswith('•'):
                    if not in_list:
                        formatted_lines.append('<ul>')
                        in_list = True
                    formatted_lines.append(f"<li>{line.strip()[1:].strip()}</li>")
                else:
                    if in_list:
                        formatted_lines.append('</ul>')
                        in_list = False
                    formatted_lines.append(f"<p>{line}</p>")
            
            if in_list:
                formatted_lines.append('</ul>')
            
            content = ''.join(formatted_lines)
        else:
            content = f"<p>{content}</p>"
        
        html += f"""
        <div class="slide">
            <h2>{slide.get('title', 'Slide')}</h2>
            {content}
        </div>
        """
    
    html += """
    </body>
    </html>
    """
    
    return html

@app.route('/api/templates', methods=['GET'])
def get_templates():
    """Get available pitch deck templates"""
    return jsonify({
        "templates": [
            {
                "id": "modern",
                "name": "Modern",
                "description": "Clean and bold design for tech startups"
            },
            {
                "id": "minimal",
                "name": "Minimal",
                "description": "Simple and elegant for any industry"
            },
            {
                "id": "startup",
                "name": "Startup",
                "description": "Vibrant and dynamic for early-stage companies"
            }
        ],
        "slide_types": [
            "title", "problem", "solution", "market", "business-model",
            "traction", "competition", "team", "financials", "ask", "custom"
        ]
    })

@app.route('/api/suggestions/investor-questions', methods=['POST'])
@limiter.limit("10 per hour")
def suggest_investor_questions():
    """Suggest potential investor questions based on pitch content"""
    try:
        data = request.json
        industry = data.get('industry', '')
        funding_stage = data.get('funding_stage', '')
        
        questions = get_investor_questions(industry, funding_stage)
        
        return jsonify({"questions": questions})
    
    except Exception as e:
        logger.error(f"Question suggestion error: {e}")
        return jsonify({"error": "Failed to generate questions"}), 500

def get_investor_questions(industry: str, funding_stage: str) -> List[str]:
    """Get common investor questions by stage and industry"""
    
    base_questions = [
        "What's your customer acquisition cost and lifetime value?",
        "How do you differentiate from competitors?",
        "What's your go-to-market strategy?",
        "How will you use the funding?",
        "What are the key risks and how will you mitigate them?",
        "What's your path to profitability?",
        "How big is the total addressable market?",
        "What's your unfair advantage?",
        "Tell me about your team's background.",
        "What are your unit economics?"
    ]
    
    stage_questions = {
        "pre-seed": [
            "Have you validated the problem with potential customers?",
            "What's your MVP timeline?",
            "How much runway will this funding provide?"
        ],
        "seed": [
            "What's your current MRR/ARR?",
            "How many customers do you have?",
            "What's your churn rate?"
        ],
        "series-a": [
            "What's your sales efficiency ratio?",
            "How will you scale the sales team?",
            "What's your expansion revenue?"
        ]
    }
    
    industry_questions = {
        "fintech": [
            "How will you handle regulatory compliance?",
            "What's your approach to security and fraud prevention?"
        ],
        "healthcare": [
            "What's your FDA approval timeline?",
            "How will you handle HIPAA compliance?"
        ],
        "saas": [
            "What's your net revenue retention?",
            "How sticky is your product?"
        ]
    }
    
    questions = base_questions.copy()
    questions.extend(stage_questions.get(funding_stage, []))
    questions.extend(industry_questions.get(industry.lower(), []))
    
    return questions[:15]  # Return top 15 questions

@app.errorhandler(429)
def ratelimit_handler(e):
    """Handle rate limit exceeded"""
    return jsonify({
        "error": "Rate limit exceeded",
        "message": str(e.description)
    }), 429

@app.errorhandler(500)
def internal_error(error):
    """Handle internal server errors"""
    logger.error(f"Internal error: {error}")
    return jsonify({
        "error": "Internal server error",
        "message": "An unexpected error occurred"
    }), 500

# Main execution
if __name__ == '__main__':
    port = int(os.getenv('PORT', 5001))
    debug = os.getenv('FLASK_ENV') == 'development'
    
    print(f"""
    🚀 Pitch Deck Builder Backend Starting...
    =========================================
    Port: {port}
    Debug: {debug}
    AI Provider: {Config.AI_PROVIDER}
    AI Available: {ai_client is not None}
    
    Endpoints:
    - GET  /                              - Frontend
    - POST /api/generate-pitch            - Generate AI pitch deck
    - POST /api/extract-pitch             - Extract from documents
    - POST /api/analyze-pitch             - Analyze pitch quality
    - POST /api/export/pptx               - Export PowerPoint
    - POST /api/export/pdf                - Export PDF
    - POST /api/suggestions/investor-questions - Get investor questions
    - GET  /api/templates                 - List templates
    - GET  /health                        - Health check
    
    To test: 
    curl -X POST http://localhost:{port}/api/generate-pitch \\
         -H "Content-Type: application/json" \\
         -d '{{"company_name": "TechCo", "problem": "Slow processes", "solution": "AI automation"}}'
    """)
    
    app.run(host='0.0.0.0', port=port, debug=debug)